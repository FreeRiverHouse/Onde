#!/usr/bin/env python3
"""
PowerPoint Social Calendar V3 - English Version with account-specific tones
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Brand colors
COLORS = {
    'magmatic': {'primary': RGBColor(138, 43, 226), 'secondary': RGBColor(180, 100, 255), 'bg': RGBColor(20, 15, 30)},
    'freeriverhouse': {'primary': RGBColor(0, 180, 200), 'secondary': RGBColor(100, 220, 230), 'bg': RGBColor(15, 30, 40)},
    'onde': {'primary': RGBColor(212, 175, 55), 'secondary': RGBColor(255, 215, 100), 'bg': RGBColor(25, 20, 15)}
}

# ============================================
# MAGMATIC POSTS - Poetic, intimate, whispered
# ============================================
MAGMATIC_POSTS = [
    # 1 - Classic from Il Cosmo Dissommerso
    """Turn of the screw, life,
shot of madness.
Three hundred sixty-five spins
around the sun,
dozens of head-spinning moments,
billions of empty loops
beside
the heart.""",

    # 2 - Reflection on pain
    """Don't rip the pain from your heart.
Don't try to put out the fire.

Soil to cultivate.
Mud to bloom from.""",

    # 3 - Inner cosmos
    """inside your head
there's the whole universe
there's a source of storms
there's the submerged cosmos""",

    # 4 - Distance and dreams (from Y Lan)
    """The boundless night separates us,
to unite we only have thoughts.
The entire night can be filled
by dreams, as long as we believe them true.""",

    # 5 - Stars
    """every atom, every element
that makes up the world
was forged in an instant
in the deep of space
by a burning star.""",

    # 6 - Bilingual wind/sea
    """The wind flows through me. Erases me.
His maid is the sea
And the soul spills her skin.""",

    # 7 - Flying
    """to dream of flying
is like setting free
an animal""",

    # 8 - Roots and stars
    """Into the earth
Towards the stars
We stretch branches
We extend hands""",

    # 9 - LA River video
    """🐸💚🌊

[sunset on the river]""",

    # 10 - Music set
    """saturday, 2:30pm
orange bridge, LA river

🎶"""
]

# ============================================
# FRH POSTS - Friendly, storytelling, real
# ============================================
FRH_POSTS = [
    # 1 - Real debugging story
    """You know that feeling when the test that's been failing for three days
suddenly passes... and you have no idea why?

Today is that day.

Celebrating by not touching anything.""",

    # 2 - Kids as QA
    """4-year-olds are the best QA testers.

They find impossible bugs.
They click where you shouldn't click.
And their feedback?

"Make it more squish."

Working on squishness.""",

    # 3 - 3 AM automation
    """The machines work while I sleep.

Test automation at 2 AM.
I wake up to green checkmarks or a to-do list.

Sometimes magic is just... delegating to robots.""",

    # 4 - Six apps, one dev
    """Six apps in parallel. One developer.

The secret? There's no secret.
Just shared architecture
and a lot of copy-paste between projects.

(copy-paste is underrated)""",

    # 5 - VR dizzy
    """VR rule number one:
if it makes you dizzy,
it'll make your users sick.

Testing on yourself isn't optional.
It's survival.""",

    # 6 - PDF margins
    """Publishing a digital book is:
- 10% writing
- 90% figuring out why the PDF margins are wrong

Every. Single. Time.""",

    # 7 - Unity 6
    """Unity 6 is fast.

Like, noticeably fast.
The editor doesn't freeze when I hit play anymore.

Small wins.""",

    # 8 - Building in public
    """Building in public means
sharing the ugly commits too.

Not every day is a breakthrough.
Sometimes it's just: "fixed typo in readme".""",

    # 9 - Polyglot dev
    """SwiftUI for the dashboard.
Unity for the games.
TypeScript for the bots.
Python for automation.

The modern indie dev is polyglot by necessity.
(and a bit by masochism)""",

    # 10 - Factory loop
    """The factory loop:
build → test → fix → repeat

The machines do the boring parts.
You sleep.
They don't.

It's a good deal."""
]

# ============================================
# ONDE POSTS - Warm, wise, Gianni style (inspired by Meditations forward)
# ============================================
ONDE_POSTS = [
    # 1 - Meditations forward style
    """Two thousand years ago, an emperor wrote notes
in a tent, at night, after days of war.

They weren't for anyone.
Just a man trying to stay sane.

And somehow, here they are. In your hands.
Free.

Sometimes the world works.""",

    # 2 - Why classics
    """You know why we publish classics?

Because a book that survived twenty centuries
has something to say.

And because wisdom shouldn't cost
more than a coffee.""",

    # 3 - Marcus quote + context
    """"You have power over your mind—not outside events.
Realize this, and you will find strength."

Marcus Aurelius wasn't a guru.
He was an emperor trying not to lose his mind
while the world crumbled around him.

Two thousand years later, the advice still holds.""",

    # 4 - Small publishers
    """We're a small publishing house in Los Angeles.
Our authors have been dead for centuries.
Our readers are everywhere.

Strange business, this one.
But we love it.""",

    # 5 - AI as tool
    """AI helps us illustrate.
Humans choose the words.

Technology amplifies culture.
It doesn't replace it.

(at least, that's the plan)""",

    # 6 - Public domain magic
    """The public domain is magic.

Shakespeare, free.
Marcus Aurelius, free.
Mary Shelley, free.

Centuries of wisdom, accessible to all.
Sometimes the internet is beautiful.""",

    # 7 - Mary Shelley
    """Mary Shelley was eighteen
when she wrote Frankenstein.

Eighteen.

In a summer of ghost stories
with Byron and Shelley,
she created the first science fiction novel.

And asked a question we still haven't answered:
what happens when we create something
we can't control?""",

    # 8 - Time
    """Books cross centuries.
We're just the bridge.

A Roman emperor speaks to you,
across two thousand years,
through a phone.

Time is strange.
Books are stranger.""",

    # 9 - Stoicism
    """Stoicism isn't about suppressing emotions.
It's about understanding them.

Marcus Aurelius knew this.
Two thousand years ago.

We're still learning.""",

    # 10 - Who we are
    """Onde.
Tech, spirituality, art.

Big words for a small operation
with big dreams.

But hey, even great empires
started with a tent in the desert."""
]

def create_title_slide(prs, title, subtitle, handle, colors):
    """Create section title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = colors['bg']
    bg.line.fill.background()

    # Accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.3), prs.slide_width, Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = colors['primary']
    bar.line.fill.background()

    # Handle
    handle_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(0.6))
    tf = handle_box.text_frame
    p = tf.paragraphs[0]
    p.text = handle
    p.font.size = Pt(28)
    p.font.color.rgb = colors['secondary']
    p.alignment = PP_ALIGN.CENTER

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = colors['primary']
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(9), Inches(0.6))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(160, 160, 160)
    p.alignment = PP_ALIGN.CENTER

def create_post_slide(prs, post_num, post_text, account, colors, tone_hint):
    """Create single post slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = colors['bg']
    bg.line.fill.background()

    # Top line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = colors['primary']
    line.line.fill.background()

    # Post number
    num_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.35), Inches(1), Inches(0.5))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"#{post_num}"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = colors['primary']

    # Account
    acc_box = slide.shapes.add_textbox(Inches(1.4), Inches(0.38), Inches(4), Inches(0.5))
    tf = acc_box.text_frame
    p = tf.paragraphs[0]
    p.text = account
    p.font.size = Pt(18)
    p.font.color.rgb = colors['secondary']

    # Tone hint (small, top right)
    tone_box = slide.shapes.add_textbox(Inches(6), Inches(0.38), Inches(3.5), Inches(0.4))
    tf = tone_box.text_frame
    p = tf.paragraphs[0]
    p.text = tone_hint
    p.font.size = Pt(12)
    p.font.italic = True
    p.font.color.rgb = RGBColor(100, 100, 100)
    p.alignment = PP_ALIGN.RIGHT

    # Main content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = post_text

    # Font size based on length
    if len(post_text) < 120:
        p.font.size = Pt(30)
    elif len(post_text) < 250:
        p.font.size = Pt(24)
    elif len(post_text) < 400:
        p.font.size = Pt(20)
    else:
        p.font.size = Pt(18)

    p.font.color.rgb = RGBColor(235, 235, 235)
    p.line_spacing = 1.5

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # === COVER ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(12, 12, 15)
    bg.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(9), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "SOCIAL CALENDAR"
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = RGBColor(212, 175, 55)
    p.alignment = PP_ALIGN.CENTER

    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1.5))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "30 Posts with Soul"
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER

    accounts_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(1))
    tf = accounts_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Onde  •  Free River House  •  Magmatic"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(140, 140, 140)
    p.alignment = PP_ALIGN.CENTER

    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = "January 2026"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(80, 80, 80)
    p.alignment = PP_ALIGN.CENTER

    # === ONDE SECTION ===
    create_title_slide(prs, "The Voice of Wisdom", "Warm • Reflective • Timeless", "@Onde_FRH", COLORS['onde'])
    for i, post in enumerate(ONDE_POSTS, 1):
        create_post_slide(prs, i, post, "@Onde_FRH", COLORS['onde'], "tone: warm, wise")

    # === FRH SECTION ===
    create_title_slide(prs, "Building in Public", "Friendly • Storytelling • Real", "@FreeRiverHouse", COLORS['freeriverhouse'])
    for i, post in enumerate(FRH_POSTS, 1):
        create_post_slide(prs, i, post, "@FreeRiverHouse", COLORS['freeriverhouse'], "tone: friendly, witty")

    # === MAGMATIC SECTION ===
    create_title_slide(prs, "Poetry & Silence", "Intimate • Poetic • Zero sales", "@magmatic__", COLORS['magmatic'])
    for i, post in enumerate(MAGMATIC_POSTS, 1):
        create_post_slide(prs, i, post, "@magmatic__", COLORS['magmatic'], "tone: whispered, poetic")

    # Save
    output = "/Users/mattiapetrucciani/Onde/Social-Calendar-30-Posts-EN.pptx"
    prs.save(output)
    print(f"✅ PowerPoint English saved: {output}")
    return output

if __name__ == "__main__":
    main()
