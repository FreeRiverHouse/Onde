#!/usr/bin/env python3
"""
PowerPoint Social Calendar V2 - Con toni specifici per account
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Colori brand
COLORS = {
    'magmatic': {'primary': RGBColor(138, 43, 226), 'secondary': RGBColor(180, 100, 255), 'bg': RGBColor(20, 15, 30)},
    'freeriverhouse': {'primary': RGBColor(0, 180, 200), 'secondary': RGBColor(100, 220, 230), 'bg': RGBColor(15, 30, 40)},
    'onde': {'primary': RGBColor(212, 175, 55), 'secondary': RGBColor(255, 215, 100), 'bg': RGBColor(25, 20, 15)}
}

# ============================================
# POST MAGMATIC - Poetico, intimo, sussurrato
# ============================================
MAGMATIC_POSTS = [
    # 1 - Classico dal Cosmo Dissommerso
    """Giro di vite, vita,
tiro di follia.
Trecentosessantacinque raggiri
intorno al sole,
dozzine di giramenti di testa,
miliardi i giri a vuoto
allato
al cuore.""",

    # 2 - Riflessione sul dolore
    """Non estirpare il male dal tuo cuore.
Non cercare di spegnere il dolore.

Terra da coltivare.
Fango su cui fiorire.""",

    # 3 - Cosmo interiore
    """nella tua testa
c'è tutto l'universo
c'è una sorgente di tempesta
c'è il cosmo dissommerso""",

    # 4 - Distanza e sogno (da Y Lan)
    """Ci separa la notte smisurata,
per unirci abbiamo solo i pensieri.
L'intera notte può essere colmata
dai sogni, finché li crediamo veri.""",

    # 5 - Stelle
    """ogni atomo, elemento
di cui è composto il mondo
lo ha forgiato in un quanto
nello spazio profondo
una stella bruciando.""",

    # 6 - Bilingue vento/mare
    """Il vento mi attraversa, mi cancella.
Il mare è la sua ancella
e l'anima ci sversa la sua scorza.

The wind flows through me. Erases me.
His maid is the sea
And the soul spills her skin.""",

    # 7 - Volare
    """sognare di volare
è come liberare
un animale""",

    # 8 - Radici e stelle
    """Dentro la terra
Verso le stelle
Tendiamo rami
Stendiamo mani""",

    # 9 - Video LA River (solo emoji)
    """🐸💚🌊

[video tramonto sul fiume]""",

    # 10 - Set musicale
    """sabato, 2:30pm
orange bridge, LA river

🎶"""
]

# ============================================
# POST FRH - Amichevole, storytelling, reale
# ============================================
FRH_POSTS = [
    # 1 - Storia vera di debugging
    """Sai quella sensazione quando il test che falliva da tre giorni
improvvisamente passa... e non sai perché?

Oggi è quel giorno.

Festeggio non toccando più niente.""",

    # 2 - Kids as QA
    """I bambini di 4 anni sono i migliori QA tester.

Trovano bug impossibili.
Cliccano dove non dovresti cliccare.
E il loro feedback?

"Fallo più squish."

Sto lavorando su squishness.""",

    # 3 - 3 AM automation
    """Le macchine lavorano mentre dormo.

Test automation alle 2 AM.
Mi sveglio con checkmark verdi o una lista di cose da fare.

A volte la magia è solo... delegare ai robot.""",

    # 4 - Sei app, un dev
    """Sei app in parallelo. Un developer.

Il segreto? Non c'è nessun segreto.
Solo architettura condivisa
e un sacco di copia-incolla tra progetti.

(il copia-incolla è sottovalutato)""",

    # 5 - VR dizzy
    """Regola VR numero uno:
se ti fa girare la testa a te,
farà vomitare gli utenti.

Testare su se stessi non è opzionale.
È sopravvivenza.""",

    # 6 - PDF margins
    """Pubblicare un libro digitale è:
- 10% scrivere
- 90% capire perché i margini del PDF sono sbagliati

Ogni. Singola. Volta.""",

    # 7 - Unity 6
    """Unity 6 è veloce.

Tipo, notevolmente veloce.
L'editor non si freeza più quando premo play.

Piccole vittorie.""",

    # 8 - Building in public
    """Building in public significa
condividere anche i commit brutti.

Non tutti i giorni sono breakthrough.
A volte è solo: "fixed typo in readme".""",

    # 9 - Polyglot dev
    """SwiftUI per la dashboard.
Unity per i giochi.
TypeScript per i bot.
Python per l'automazione.

L'indie dev moderno è poliglotta per necessità.
(e un po' per masochismo)""",

    # 10 - Factory loop
    """Il loop della fabbrica:
build → test → fix → repeat

Le macchine fanno le parti noiose.
Tu dormi.
Loro no.

È un buon accordo."""
]

# ============================================
# POST ONDE - Caldo, saggio, stile Gianni (ispirato a Meditations forward)
# ============================================
ONDE_POSTS = [
    # 1 - Stile forward Meditations
    """Duemila anni fa, un imperatore scriveva appunti
in una tenda, di notte, dopo giorni di guerra.

Non erano per nessuno.
Solo un uomo che cercava di restare sano di mente.

E in qualche modo, eccoli qui. Nelle tue mani.
Gratis.

A volte il mondo funziona.""",

    # 2 - Perché classici
    """Sai perché pubblichiamo classici?

Perché un libro che ha attraversato venti secoli
ha qualcosa da dire.

E perché la saggezza non dovrebbe costare
più di un caffè.""",

    # 3 - Citazione Marcus + contesto
    """"Hai potere sulla tua mente—non sugli eventi esterni.
Realizza questo, e troverai forza."

Marco Aurelio non era un guru.
Era un imperatore che cercava di non impazzire
mentre il mondo crollava intorno a lui.

Duemila anni dopo, il consiglio regge ancora.""",

    # 4 - Piccoli editori
    """Siamo una piccola casa editrice a Los Angeles.
I nostri autori sono morti da secoli.
I nostri lettori sono ovunque.

Strano mestiere, questo.
Ma ci piace.""",

    # 5 - AI come strumento
    """L'AI ci aiuta a illustrare.
Gli umani scelgono le parole.

La tecnologia amplifica la cultura.
Non la sostituisce.

(almeno, questo è il piano)""",

    # 6 - Public domain magic
    """Il public domain è magia.

Shakespeare, gratis.
Marco Aurelio, gratis.
Mary Shelley, gratis.

Secoli di saggezza, accessibili a tutti.
Internet a volte è bellissimo.""",

    # 7 - Mary Shelley
    """Mary Shelley aveva diciotto anni
quando scrisse Frankenstein.

Diciotto.

In un'estate di storie di fantasmi
con Byron e Shelley,
creò il primo romanzo di fantascienza.

E fece una domanda che ancora non abbiamo risposto:
cosa succede quando creiamo qualcosa
che non possiamo controllare?""",

    # 8 - Il tempo
    """I libri attraversano i secoli.
Noi siamo solo il ponte.

Un imperatore romano parla a te,
attraverso duemila anni,
attraverso un telefono.

Il tempo è strano.
I libri lo sono di più.""",

    # 9 - Stoicismo
    """Lo stoicismo non è sopprimere le emozioni.
È capirle.

Marco Aurelio lo sapeva.
Duemila anni fa.

Noi stiamo ancora imparando.""",

    # 10 - Chi siamo
    """Onde.
Tech, spiritualità, arte.

Parole grosse per un'operazione piccola
con sogni grandi.

Ma ehi, anche i grandi imperi
sono iniziati con una tenda nel deserto."""
]

def create_title_slide(prs, title, subtitle, handle, colors):
    """Crea slide titolo sezione"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = colors['bg']
    bg.line.fill.background()

    # Accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.3), prs.slide_width, Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = colors['primary']
    bar.line.fill.background()

    # Handle
    handle_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(0.6))
    tf = handle_box.text_frame
    p = tf.paragraphs[0]
    p.text = handle
    p.font.size = Pt(28)
    p.font.color.rgb = colors['secondary']
    p.alignment = PP_ALIGN.CENTER

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = colors['primary']
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(9), Inches(0.6))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(160, 160, 160)
    p.alignment = PP_ALIGN.CENTER

def create_post_slide(prs, post_num, post_text, account, colors, tone_hint):
    """Crea slide singolo post"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = colors['bg']
    bg.line.fill.background()

    # Top line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = colors['primary']
    line.line.fill.background()

    # Post number
    num_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.35), Inches(1), Inches(0.5))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"#{post_num}"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = colors['primary']

    # Account
    acc_box = slide.shapes.add_textbox(Inches(1.4), Inches(0.38), Inches(4), Inches(0.5))
    tf = acc_box.text_frame
    p = tf.paragraphs[0]
    p.text = account
    p.font.size = Pt(18)
    p.font.color.rgb = colors['secondary']

    # Tone hint (piccolo, in alto a destra)
    tone_box = slide.shapes.add_textbox(Inches(6), Inches(0.38), Inches(3.5), Inches(0.4))
    tf = tone_box.text_frame
    p = tf.paragraphs[0]
    p.text = tone_hint
    p.font.size = Pt(12)
    p.font.italic = True
    p.font.color.rgb = RGBColor(100, 100, 100)
    p.alignment = PP_ALIGN.RIGHT

    # Main content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = post_text

    # Font size based on length
    if len(post_text) < 120:
        p.font.size = Pt(30)
    elif len(post_text) < 250:
        p.font.size = Pt(24)
    elif len(post_text) < 400:
        p.font.size = Pt(20)
    else:
        p.font.size = Pt(18)

    p.font.color.rgb = RGBColor(235, 235, 235)
    p.line_spacing = 1.5

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # === COVER ===
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(12, 12, 15)
    bg.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(9), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "SOCIAL CALENDAR"
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = RGBColor(212, 175, 55)
    p.alignment = PP_ALIGN.CENTER

    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1.5))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "30 Post con Anima"
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER

    accounts_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(1))
    tf = accounts_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Onde  •  Free River House  •  Magmatic"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(140, 140, 140)
    p.alignment = PP_ALIGN.CENTER

    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Gennaio 2026"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(80, 80, 80)
    p.alignment = PP_ALIGN.CENTER

    # === ONDE SECTION ===
    create_title_slide(prs, "La Voce della Saggezza", "Caldo • Riflessivo • Stile Gianni", "@Onde_FRH", COLORS['onde'])
    for i, post in enumerate(ONDE_POSTS, 1):
        create_post_slide(prs, i, post, "@Onde_FRH", COLORS['onde'], "tono: caldo, saggio")

    # === FRH SECTION ===
    create_title_slide(prs, "Building in Public", "Amichevole • Storytelling • Reale", "@FreeRiverHouse", COLORS['freeriverhouse'])
    for i, post in enumerate(FRH_POSTS, 1):
        create_post_slide(prs, i, post, "@FreeRiverHouse", COLORS['freeriverhouse'], "tono: amichevole, ironico")

    # === MAGMATIC SECTION ===
    create_title_slide(prs, "Poesia & Silenzio", "Intimo • Poetico • Zero vendita", "@magmatic__", COLORS['magmatic'])
    for i, post in enumerate(MAGMATIC_POSTS, 1):
        create_post_slide(prs, i, post, "@magmatic__", COLORS['magmatic'], "tono: sussurrato, poetico")

    # Save
    output = "/Users/mattiapetrucciani/Onde/Social-Calendar-30-Posts-V2.pptx"
    prs.save(output)
    print(f"✅ PowerPoint V2 salvato: {output}")
    return output

if __name__ == "__main__":
    main()
