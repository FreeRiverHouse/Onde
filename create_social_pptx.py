#!/usr/bin/env python3
"""
Crea PowerPoint lussuoso con 30 post social (10 per account)
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.util import Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor as RgbColor
import os

# Colori brand
COLORS = {
    'magmatic': {'primary': RgbColor(75, 0, 130), 'secondary': RgbColor(138, 43, 226), 'bg': RgbColor(25, 25, 35)},
    'freeriverhouse': {'primary': RgbColor(0, 100, 150), 'secondary': RgbColor(0, 180, 200), 'bg': RgbColor(20, 35, 45)},
    'onde': {'primary': RgbColor(212, 175, 55), 'secondary': RgbColor(180, 140, 40), 'bg': RgbColor(30, 25, 20)}
}

# Post per Magmatic (poesia)
MAGMATIC_POSTS = [
    """Giro di vite, vita,
tiro di follia.
Trecentosessantacinque raggiri
intorno al sole,
dozzine di giramenti di testa,
miliardi i giri a vuoto
allato
al cuore.""",
    """Non estirpare il male dal tuo cuore.
Non cercare di spegnere il dolore.

Terra da coltivare.
Fango su cui fiorire.""",
    """nella tua testa
c'è tutto l'universo
c'è una sorgente di tempesta
c'è il cosmo dissommerso""",
    """Ci separa la notte smisurata,
per unirci abbiamo solo i pensieri.
L'intera notte può essere colmata
dai sogni, finché li crediamo veri.""",
    """ogni atomo, elemento
di cui è composto il mondo
lo ha forgiato in un quanto
nello spazio profondo
una stella bruciando.""",
    """Il vento mi attraversa, mi cancella.
Il mare è la sua ancella
e l'anima ci sversa la sua scorza.

The wind flows through me. Erases me.
His maid is the sea
And the soul spills her skin.""",
    """sognare di volare
è come liberare
un animale""",
    """Dentro la terra
Verso le stelle
Tendiamo rami
Stendiamo mani""",
    """Radici nel buio,
foglie verso il sole.
Ogni giorno un'alba,
ogni notte una scuola.""",
    """Siamo polvere di stelle
che ha imparato a guardarsi allo specchio.
E a chiedersi: perché?"""
]

# Post per FreeRiverHouse (tech/building in public)
FRH_POSTS = [
    "Building apps for kids is different. Every button needs to work on the first tap. Every error needs a friendly message. No second chances.",
    "Our Unity test automation runs at 2 AM. I wake up to either green checkmarks or a to-do list. The machines never sleep.",
    "Six apps in parallel, one developer. The secret? Shared architecture and a lot of copy-paste between projects.",
    "VR development rule: if it makes you dizzy, it'll make your users sick. Testing on yourself is non-negotiable.",
    "The factory loop: build, test, fix, repeat. Let the machines do the boring parts while you sleep.",
    "Kids don't read tutorials. They tap everything. Design for chaos, build for curiosity.",
    "Publishing a digital book is 10% writing, 90% figuring out why the PDF margins are wrong. Every. Single. Time.",
    "Unity 6 is fast. Like, noticeably fast. The editor doesn't freeze when I hit play anymore. Small wins.",
    "Building in public means sharing the ugly commits too. Not every day is a breakthrough.",
    "SwiftUI for the dashboard, Unity for the games, TypeScript for the bots. Modern indie dev is polyglot by necessity."
]

# Post per Onde (publishing/philosophy)
ONDE_POSTS = [
    "Two thousand years ago, an emperor wrote notes to himself about staying calm. Today, you can read them for free. That's the magic of public domain.",
    "We don't just publish books. We bring ancient wisdom into the digital age. Illustrated. Beautiful. Accessible.",
    "Marcus Aurelius didn't write for fame. He wrote to survive. The Meditations were never meant to be published. Now they're one of humanity's greatest gifts.",
    "AI helps us illustrate. Humans choose the words. Technology amplifies culture, it doesn't replace it.",
    "Every book we publish is free. Not because it has no value. Because wisdom shouldn't have a price tag.",
    "Stoicism isn't about suppressing emotions. It's about understanding them. Marcus Aurelius knew this 2000 years ago.",
    "The best books aren't written. They're discovered. In attics. In monasteries. In the public domain.",
    "We're a small publisher in Los Angeles. Our authors have been dead for centuries. Our readers are everywhere.",
    "Tech + Spirituality + Art. That's Onde. Classical wisdom, modern design, universal access.",
    "\"You have power over your mind—not outside events. Realize this, and you will find strength.\" — Marcus Aurelius"
]

def create_title_slide(prs, title, subtitle, colors):
    """Crea slide titolo sezione"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = colors['bg']
    background.line.fill.background()

    # Accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.2), prs.slide_width, Inches(0.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = colors['primary']
    bar.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = colors['primary']
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RgbColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER

def create_post_slide(prs, post_num, total_in_section, post_text, account, colors):
    """Crea slide singolo post"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Background gradient effect (solid for simplicity)
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = colors['bg']
    background.line.fill.background()

    # Top accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
    line.fill.solid()
    line.fill.fore_color.rgb = colors['primary']
    line.line.fill.background()

    # Post number badge
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(0.4), Inches(1.2), Inches(0.6))
    badge.fill.solid()
    badge.fill.fore_color.rgb = colors['primary']
    badge.line.fill.background()

    badge_text = badge.text_frame
    badge_text.paragraphs[0].text = f"POST {post_num}"
    badge_text.paragraphs[0].font.size = Pt(16)
    badge_text.paragraphs[0].font.bold = True
    badge_text.paragraphs[0].font.color.rgb = RgbColor(255, 255, 255)
    badge_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    badge_text.word_wrap = False

    # Account handle
    handle_box = slide.shapes.add_textbox(Inches(1.8), Inches(0.45), Inches(5), Inches(0.5))
    tf = handle_box.text_frame
    p = tf.paragraphs[0]
    p.text = account
    p.font.size = Pt(18)
    p.font.color.rgb = colors['secondary']

    # Main post content
    content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(8.8), Inches(4.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = post_text

    # Adjust font size based on content length
    if len(post_text) < 150:
        p.font.size = Pt(32)
    elif len(post_text) < 300:
        p.font.size = Pt(26)
    else:
        p.font.size = Pt(22)

    p.font.color.rgb = RgbColor(240, 240, 240)
    p.line_spacing = 1.4

    # Progress indicator
    progress_box = slide.shapes.add_textbox(Inches(8), Inches(6.8), Inches(1.5), Inches(0.3))
    tf = progress_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{post_num}/{total_in_section}"
    p.font.size = Pt(14)
    p.font.color.rgb = RgbColor(100, 100, 100)
    p.alignment = PP_ALIGN.RIGHT

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Cover slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RgbColor(15, 15, 20)
    bg.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "SOCIAL CALENDAR"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RgbColor(212, 175, 55)
    p.alignment = PP_ALIGN.CENTER

    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "30 Post di Qualità\nOnde • Free River House • Magmatic"
    p.font.size = Pt(28)
    p.font.color.rgb = RgbColor(180, 180, 180)
    p.alignment = PP_ALIGN.CENTER

    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Gennaio 2026"
    p.font.size = Pt(18)
    p.font.color.rgb = RgbColor(100, 100, 100)
    p.alignment = PP_ALIGN.CENTER

    # ONDE section
    create_title_slide(prs, "@ONDE_FRH", "Publishing • Philosophy • Wisdom", COLORS['onde'])
    for i, post in enumerate(ONDE_POSTS, 1):
        create_post_slide(prs, i, 10, post, "@Onde_FRH", COLORS['onde'])

    # FRH section
    create_title_slide(prs, "@FREERIVERHOUSE", "Building in Public • Tech • Indie Dev", COLORS['freeriverhouse'])
    for i, post in enumerate(FRH_POSTS, 1):
        create_post_slide(prs, i, 10, post, "@FreeRiverHouse", COLORS['freeriverhouse'])

    # MAGMATIC section
    create_title_slide(prs, "@MAGMATIC__", "Poetry • Art • Soul", COLORS['magmatic'])
    for i, post in enumerate(MAGMATIC_POSTS, 1):
        create_post_slide(prs, i, 10, post, "@magmatic__", COLORS['magmatic'])

    # Save
    output_path = "/Users/mattiapetrucciani/Onde/Social-Calendar-30-Posts.pptx"
    prs.save(output_path)
    print(f"PowerPoint salvato: {output_path}")
    return output_path

if __name__ == "__main__":
    main()
