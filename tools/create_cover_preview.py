#!/usr/bin/env python3
"""
Creates a PowerPoint preview of album covers for Onde Lounge
Since we can't download directly from Grok, this creates placeholder slides
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Track names
TRACKS = [
    "1. Midnight Windows",
    "2. Neon Rain",
    "3. Desert Between The Stars",
    "4. Joshua Tree, 4 AM",
    "5. Golden Hour Static",
    "6. River Sunset",
    "7. Velvet Static",
    "8. 2 AM, Onde",
    "9. Slow Motion Sunset",
    "10. Prima Onda",
    "11. Velvet Hours"
]

# Cover descriptions (for visual reference)
COVER_DESCRIPTIONS = [
    "City lights through rain - blue/purple gradient with geometric windows",
    "Wet streets reflection - neon pink and teal gradient",
    "Vast night sky - deep blue with golden stars",
    "Desert starscape - warm orange horizon fading to purple sky",
    "Warm sunset gradient - golden orange to soft pink",
    "Golden water reflection - river with sunset colors",
    "Deep purple texture - velvet gradient with soft particles",
    "Ocean waves at night - deep blue with moonlight reflection",
    "Orange pink gradient - smooth sunset colors",
    "Gentle morning wave - soft teal and gold",
    "Soft piano keys silhouette - dark blue with golden accents"
]

def create_pptx():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Dark background
    from pptx.enum.shapes import MSO_SHAPE
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(15, 15, 25)
    bg.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "ONDE LOUNGE"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(212, 175, 55)  # Gold
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "11 Album Covers - Electronic Chill Collection"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(150, 150, 180)
    p.alignment = PP_ALIGN.CENTER

    # Create slide for each track
    for i, (track, desc) in enumerate(zip(TRACKS, COVER_DESCRIPTIONS)):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Dark background with gradient effect simulation
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        # Alternate between dark blue and dark purple backgrounds
        if i % 2 == 0:
            bg.fill.fore_color.rgb = RGBColor(15, 20, 35)
        else:
            bg.fill.fore_color.rgb = RGBColor(25, 15, 35)
        bg.line.fill.background()

        # Track number with gold accent
        num_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(1), Inches(0.8))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"#{i+1}"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(212, 175, 55)

        # Track name
        name_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(1))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        # Remove number from track name for display
        track_name = track.split(". ", 1)[1] if ". " in track else track
        p.text = track_name
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        # Cover description box (simulates the cover)
        cover_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                           Inches(2.5), Inches(2.8),
                                           Inches(5), Inches(3.5))
        cover_box.fill.solid()
        cover_box.fill.fore_color.rgb = RGBColor(30, 40, 60)
        cover_box.line.color.rgb = RGBColor(212, 175, 55)
        cover_box.line.width = Pt(2)

        # Description inside cover box
        desc_box = slide.shapes.add_textbox(Inches(2.7), Inches(3.5), Inches(4.6), Inches(2))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"Cover Style:\n{desc}"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(180, 180, 200)
        p.alignment = PP_ALIGN.CENTER

        # Footer with Onde branding
        footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(9), Inches(0.5))
        tf = footer_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Onde Lounge • Electronic Chill Collection • 2026"
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(100, 100, 120)
        p.alignment = PP_ALIGN.CENTER

    # Save
    output = "/Users/mattiapetrucciani/Onde/Onde-Lounge-Covers-Preview.pptx"
    prs.save(output)
    print(f"✅ PowerPoint saved: {output}")
    return output

if __name__ == "__main__":
    create_pptx()
